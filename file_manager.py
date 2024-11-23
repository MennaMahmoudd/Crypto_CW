import tkinter as tk
from tkinter import messagebox, filedialog
import hashlib
import os
import re
from Crypto.Cipher import AES
from Crypto.Util.Padding import pad
from Crypto.Util.Padding import unpad
from Crypto.Random import get_random_bytes
from docx import Document 
from PyPDF2 import PdfReader, PdfWriter
from pptx import Presentation
import openpyxl


USER_DATA_FILE = "users.txt"
UPLOADS_DIR = "uploads"  # Directory to store uploaded files


# Ensure the uploads directory exists
if not os.path.exists(UPLOADS_DIR):
    os.makedirs(UPLOADS_DIR)

# Global variable to keep track of UI components
current_widgets = []

# Function to clear all widgets in the current UI
def clear_ui():
    for widget in current_widgets:
        widget.pack_forget()
    current_widgets.clear()

# Function to hash passwords
def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()

# Function to save user data
def save_user(username, hashed_password):
    with open(USER_DATA_FILE, "a") as file:
        file.write(f"{username},{hashed_password}\n")

# Function to verify username and password
def verify_user(username, password):
    if not os.path.exists(USER_DATA_FILE):
        return False  # No users file yet
    hashed_input_password = hash_password(password)
    with open(USER_DATA_FILE, "r") as file:
        for line in file:
            saved_username, saved_password = line.strip().split(",")
            if username == saved_username and hashed_input_password == saved_password:
                return True
    return False

# Function to validate password complexity
def is_password_complex(password):
    if len(password) < 8:
        return False
    if not re.search(r"[A-Z]", password):  # At least one uppercase letter
        return False
    if not re.search(r"[a-z]", password):  # At least one lowercase letter
        return False
    if not re.search(r"\d", password):  # At least one digit
        return False
    if not re.search(r"[!@#$%^&*(),.?\":{}|<>_]", password):  # At least one special character
        return False
    return True

# Signup function
def signup():
    username = entry_username.get()
    password = entry_password.get()

    if not username or not password:
        messagebox.showwarning("Input Error", "Both fields are required!")
        return

    if not is_password_complex(password):
        messagebox.showwarning(
            "Weak Password",
            "Password must be at least 8 characters long, contain uppercase and lowercase letters, a number, and a special character.",
        )
        return

    hashed_password = hash_password(password)
    save_user(username, hashed_password)
    messagebox.showinfo("Signup Success", "Account created successfully! Please log in now.")
    
     # Clear the username and password fields
    entry_username.delete(0, tk.END)
    entry_password.delete(0, tk.END)
    show_login_screen()

# Login function
def login():
    username = entry_username.get()
    password = entry_password.get()

    if not username or not password:
        messagebox.showwarning("Input Error", "Both fields are required!")
        return

    if verify_user(username, password):
        messagebox.showinfo("Login Success", f"Welcome, {username}!")
        show_upload_screen()  # Show the upload screen on successful login
    else:
        messagebox.showerror("Login Failed", "Invalid username or password! Please sign up if you don't have an account.")
        entry_password.delete(0, tk.END)

# Function to show login screen
def show_login_screen():
    clear_ui()
    label_title.config(text="File Manager - Please Sign Up or Log In", font=("Helvetica", 18))
    label_title.pack(pady=10)
    current_widgets.append(label_title)

    label_username.pack()
    current_widgets.append(label_username)

    entry_username.pack(pady=5)
    current_widgets.append(entry_username)

    label_password.pack()
    current_widgets.append(label_password)

    entry_password.pack(pady=5)
    current_widgets.append(entry_password)

    button_login.pack(pady=10)
    current_widgets.append(button_login)

    button_signup.pack(pady=10)
    current_widgets.append(button_signup)
    
# Function to show file upload screen
def show_upload_screen():
    clear_ui()
    label_title.config(text="File Upload", font=("Helvetica", 18))
    label_title.pack(pady=10)
    current_widgets.append(label_title)

    label_info = tk.Label(root, text="Select a file to upload", font=("Helvetica", 12))
    label_info.pack(pady=10)
    current_widgets.append(label_info)

    button_browse = tk.Button(root, text="upload file", command=upload_file, font=("Helvetica", 12))
    button_browse.pack(pady=10)
    current_widgets.append(button_browse)

    button_logout = tk.Button(root, text="Logout", command=show_login_screen, font=("Helvetica", 12))
    button_logout.pack(pady=10)
    current_widgets.append(button_logout)
    
    

def encrypt_file(file_path, destination_path):
    key = get_random_bytes(16)  # AES requires a 16-byte key
    cipher = AES.new(key, AES.MODE_CBC)  # Using CBC mode for AES encryption
    with open(file_path, 'rb') as file:
        file_data = file.read()
        encrypted_data = cipher.encrypt(pad(file_data, AES.block_size))

    with open(destination_path, 'wb') as enc_file:
        # Write the IV and the encrypted data
        enc_file.write(cipher.iv)
        enc_file.write(encrypted_data)

    # Save the key in a separate file (or database, etc.)
    key_file_path = destination_path + ".key"
    with open(key_file_path, 'wb') as key_file:
        key_file.write(key)

    return key_file_path  # Return the path to the key file


# Function to handle file upload
def upload_file():
    file_path = filedialog.askopenfilename(title="Select a file", filetypes=(("Word Files", "*.docx"), ("PDF Files", "*.pdf"), ("PowerPoint Files", "*.pptx"), ("Excel Files", "*.xlsx"), ("All Files", "*.*")))
    if file_path:
        file_name = os.path.basename(file_path)
        destination = os.path.join(UPLOADS_DIR, file_name)

        if os.path.exists(destination):
            messagebox.showwarning("File Already Uploaded", f"The file '{file_name}' has already been uploaded.")
            return
        

        with open(file_path, 'rb') as source_file:
            with open(destination, 'wb') as dest_file:
                dest_file.write(source_file.read())
                
                
        encryption_key = encrypt_file(file_path, destination)
        messagebox.showinfo("File Upload", f"File '{file_name}' uploaded successfully!")
        
# Function to show uploaded files
def show_uploaded_files():
    clear_ui()
    label_title.config(text="Uploaded Files", font=("Helvetica", 18))
    label_title.pack(pady=10)
    current_widgets.append(label_title)

    # Get the list of uploaded files
    uploaded_files = os.listdir(UPLOADS_DIR)
    if not uploaded_files:
        messagebox.showinfo("No Files", "No files have been uploaded yet.")
        show_upload_screen()
        return

    # Display uploaded files
    for file in uploaded_files:
        file_label = tk.Label(root, text=file, font=("Helvetica", 12))
        file_label.pack(pady=5)
        current_widgets.append(file_label)
        
      # Display uploaded files as clickable buttons
    for file in uploaded_files:
        file_button = tk.Button(root, text=file, font=("Helvetica", 12), command=lambda file=file: file_options(file))
        file_button.pack(pady=5)
        current_widgets.append(file_button)
        


    # Back button to go back to the upload screen
    button_back = tk.Button(root, text="Back", command=show_upload_screen, font=("Helvetica", 12))
    button_back.pack(pady=10)
    current_widgets.append(button_back)
    
# Function to handle the options for a selected file
def file_options(file_name):
    clear_ui()

    label_title.config(text=f"Options for {file_name}", font=("Helvetica", 18))
    label_title.pack(pady=10)
    current_widgets.append(label_title)

    # Button to delete the file
    button_delete = tk.Button(root, text="Delete File", command=lambda: delete_file(file_name), font=("Helvetica", 12))
    button_delete.pack(pady=10)
    current_widgets.append(button_delete)

    # Button to rename the file
    button_rename = tk.Button(root, text="Rename File", command=lambda: rename_file(file_name), font=("Helvetica", 12))
    button_rename.pack(pady=10)
    current_widgets.append(button_rename)
    
      # Button to edit the file (added here)
    button_edit = tk.Button(root, text="Edit File", command=lambda: edit_file_in_window(file_name), font=("Helvetica", 12))
    button_edit.pack(pady=10)
    current_widgets.append(button_edit)
    
     # Button to decrypt and download the file
    button_download_decrypt = tk.Button(root, text="Download File", command=lambda: download_decrypt_file(file_name), font=("Helvetica", 12))
    button_download_decrypt.pack(pady=10)
    current_widgets.append(button_download_decrypt)
    

    # Back button to go back to the uploaded files list
    button_back = tk.Button(root, text="Back", command=show_uploaded_files, font=("Helvetica", 12))
    button_back.pack(pady=10)
    current_widgets.append(button_back)

# Function to delete a selected file
def delete_file(file_name):
    file_path = os.path.join(UPLOADS_DIR, file_name)
    try:
        os.remove(file_path)
        messagebox.showinfo("File Deleted", f"The file '{file_name}' has been deleted successfully!")
        show_uploaded_files()  # Refresh the file list
    except Exception as e:
        messagebox.showerror("Error", f"An error occurred while deleting the file: {str(e)}")

# Function to rename a selected file
def rename_file(file_name):
    def apply_rename():
        new_name = entry_new_name.get()
        if not new_name:
            messagebox.showwarning("Invalid Name", "Please enter a valid name.")
            return
        
        old_path = os.path.join(UPLOADS_DIR, file_name)
        new_path = os.path.join(UPLOADS_DIR, new_name)

        if os.path.exists(new_path):
            messagebox.showwarning("File Exists", f"A file with the name '{new_name}' already exists.")
            return

        try:
            os.rename(old_path, new_path)
            messagebox.showinfo("File Renamed", f"The file has been renamed to '{new_name}' successfully!")
            show_uploaded_files()  # Refresh the file list
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred while renaming the file: {str(e)}")
    
    clear_ui()

    # Prompt for a new file name
    label_title.config(text=f"Rename {file_name}", font=("Helvetica", 18))
    label_title.pack(pady=10)
    current_widgets.append(label_title)

    label_new_name = tk.Label(root, text="Enter new file name:", font=("Helvetica", 12))
    label_new_name.pack(pady=10)
    current_widgets.append(label_new_name)

    entry_new_name = tk.Entry(root, font=("Helvetica", 12))
    entry_new_name.pack(pady=10)
    entry_new_name.insert(0, file_name)  # Set the current file name as the default
    current_widgets.append(entry_new_name)

    button_apply_rename = tk.Button(root, text="Apply Rename", command=apply_rename, font=("Helvetica", 12))
    button_apply_rename.pack(pady=10)
    current_widgets.append(button_apply_rename)

    button_cancel = tk.Button(root, text="Cancel", command=show_uploaded_files, font=("Helvetica", 12))
    button_cancel.pack(pady=10)
    current_widgets.append(button_cancel)
    
# Function to edit a file in a new window
def edit_file_in_window(file_name):
    label_title.config(text=f"Edit File: {file_name}", font=("Helvetica", 18))

    file_path = os.path.join(UPLOADS_DIR, file_name)

    # Check if the file is a .docx, .pdf, .pptx, or .xlsx file and handle accordingly
    if file_name.endswith(".docx"):
        try:
            doc = Document(file_path)
            content = "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            messagebox.showerror("Error", f"Unable to open DOCX file: {e}")
            show_uploaded_files()
            return
    elif file_name.endswith(".pdf"):
        try:
            pdf_reader = PdfReader(file_path)
            content = "\n".join([page.extract_text() for page in pdf_reader.pages])
        except Exception as e:
            messagebox.showerror("Error", f"Unable to open PDF file: {e}")
            show_uploaded_files()
            return
    elif file_name.endswith(".pptx"):
        try:
            prs = Presentation(file_path)
            content = "\n".join([slide.shapes.title.text if slide.shapes.title else '' for slide in prs.slides])
        except Exception as e:
            messagebox.showerror("Error", f"Unable to open PowerPoint file: {e}")
            show_uploaded_files()
            return
    elif file_name.endswith(".xlsx"):
        try:
            wb = openpyxl.load_workbook(file_path)
            sheet = wb.active
            content = "\n".join([str(cell.value) for row in sheet.iter_rows() for cell in row])
        except Exception as e:
            messagebox.showerror("Error", f"Unable to open Excel file: {e}")
            show_uploaded_files()
            return
    else:
        content = ""
   
   
       # Create a Text widget to display the file content
    text_edit = tk.Text(root, wrap=tk.WORD, height=15, width=40)
    text_edit.insert(tk.END, content)
    text_edit.pack(pady=10)

    def save_edits():
        new_content = text_edit.get("1.0", tk.END).strip()
        if file_name.endswith(".docx"):
            try:
                doc = Document(file_path)
                for para in doc.paragraphs:
                    para.clear()  # Clear existing content
                doc.add_paragraph(new_content)
                doc.save(file_path)
                messagebox.showinfo("Edit Success", f"File '{file_name}' edited successfully!")
                messagebox.showinfo("Program Close", "File successfully saved and program will now close.")
                root.quit()  # This will close the Tkinter application
            except Exception as e:
                messagebox.showerror("Error", f"Failed to save DOCX edits: {e}")
        elif file_name.endswith(".pdf"):
            try:
                pdf_writer = PdfWriter()
                pdf_reader = PdfReader(file_path)
                for page in pdf_reader.pages:
                    pdf_writer.add_page(page)
                # Add new content to the first page (example)
                pdf_writer.pages[0].extract_text()  # You need to modify content here for editing
                with open(file_path, 'wb') as f:
                    pdf_writer.write(f)
                messagebox.showinfo("Edit Success", f"File '{file_name}' edited successfully!")
                messagebox.showinfo("Program Close", "File successfully saved and program will now close.")
                root.quit()  # This will close the Tkinter application
            except Exception as e:
                messagebox.showerror("Error", f"Failed to save PDF edits: {e}")
        elif file_name.endswith(".pptx"):
            try:
                prs = Presentation(file_path)
                # Update presentation content with new content
                prs.slides[0].shapes.title.text = new_content  # Modify as needed
                prs.save(file_path)
                messagebox.showinfo("Edit Success", f"File '{file_name}' edited successfully!")
                messagebox.showinfo("Program Close", "File successfully saved and program will now close.")
                root.quit()  # This will close the Tkinter application
            except Exception as e:
                messagebox.showerror("Error", f"Failed to save PowerPoint edits: {e}")
        elif file_name.endswith(".xlsx"):
            try:
                wb = openpyxl.load_workbook(file_path)
                sheet = wb.active
                # Update Excel content with new content (simplified)
                for row in sheet.iter_rows():
                    for cell in row:
                        cell.value = new_content  # Update cells with new content
                wb.save(file_path)
                messagebox.showinfo("Edit Success", f"File '{file_name}' edited successfully!")
                      # Close the program after saving
                messagebox.showinfo("Program Close", "File successfully saved and program will now close.")
                root.quit()  # This will close the Tkinter application
            except Exception as e:
                messagebox.showerror("Error", f"Failed to save Excel edits: {e}")


    # Save button to save the edits
    button_save = tk.Button(root, text="Save Edits", command=save_edits, font=("Helvetica", 12))
    button_save.pack(pady=10)

    # Back button to go back to the uploaded files list
    button_back = tk.Button(root, text="Back", command=show_uploaded_files, font=("Helvetica", 12))
    button_back.pack(pady=10)


def decrypt_file(encrypted_file_path, key_file_path, destination_path):
    try:
        with open(encrypted_file_path, 'rb') as enc_file:
            iv = enc_file.read(16)  # AES IV is 16 bytes
            encrypted_data = enc_file.read()

        # Retrieve the key
        with open(key_file_path, 'rb') as key_file:
            key = key_file.read()

        # Use the correct IV and key for decryption
        cipher = AES.new(key, AES.MODE_CBC, iv)
        
        # Unpad the decrypted data
        decrypted_data = unpad(cipher.decrypt(encrypted_data), AES.block_size)

        # Save the decrypted file
        with open(destination_path, 'wb') as dec_file:
            dec_file.write(decrypted_data)

        return True
    except Exception as e:
        messagebox.showerror("Error", f"An error occurred while decrypting the file: {e}")
        return False

   


# Function to handle decryption and download of the file
def download_decrypt_file(file_name):
    encrypted_file_path = os.path.join(UPLOADS_DIR, file_name)
    key_file_path = encrypted_file_path + ".key"  # Assuming the key is stored with this extension
    
    # Check if the encrypted file and key file exist
    if not os.path.exists(encrypted_file_path):
        messagebox.showerror("Error", "Encrypted file not found.")
        return
    if not os.path.exists(key_file_path):
        messagebox.showerror("Error", "Key file not found.")
        return

    # Ask the user where to save the decrypted file
    decrypted_file_path = filedialog.asksaveasfilename(
        defaultextension=".dec", filetypes=[("All Files", "*.*")], title="Save Decrypted File"
    )

    if not decrypted_file_path:
        return  # User canceled the save dialog

    # Perform decryption
    if decrypt_file(encrypted_file_path, key_file_path, decrypted_file_path):
        messagebox.showinfo("Download Success", f"File successfully and saved as {decrypted_file_path}.")
    else:
        messagebox.showerror("Error", "Decryption failed.")

    # Ask user for a destination to save the decrypted file
    destination_path = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("DOCX Files", "*.docx"), ("All Files", "*.*")])

    if destination_path:
        success = decrypt_file(encrypted_file_path, key_file_path, destination_path)
        if success:
            messagebox.showinfo("The file is Successfully", f"File  and saved to '{destination_path}'.")
        else:
            messagebox.showerror("Decryption Failed", "Failed to decrypt the file.")





# Function to show file upload screen with options
def show_upload_screen():
    clear_ui()
    label_title.config(text="File Upload", font=("Helvetica", 18))
    label_title.pack(pady=10)
    current_widgets.append(label_title)

    # Option to upload a file
    button_upload = tk.Button(root, text="Upload a File", command=upload_file, font=("Helvetica", 12))
    button_upload.pack(pady=10)
    current_widgets.append(button_upload)

    # Option to show uploaded files
    button_show_files = tk.Button(root, text="Show Uploaded Files", command=show_uploaded_files, font=("Helvetica", 12))
    button_show_files.pack(pady=10)
    current_widgets.append(button_show_files)

    # Logout button
    button_logout = tk.Button(root, text="Logout", command=show_login_screen, font=("Helvetica", 12))
    button_logout.pack(pady=10)
    current_widgets.append(button_logout)

        
# Create the main window
root = tk.Tk()
root.title("File Manager")

# Create shared widgets
label_title = tk.Label(root, text="", font=("Helvetica", 18))

label_username = tk.Label(root, text="Username", font=("Helvetica", 12))
entry_username = tk.Entry(root, font=("Helvetica", 12))

label_password = tk.Label(root, text="Password", font=("Helvetica", 12))
entry_password = tk.Entry(root, font=("Helvetica", 12), show="*")

button_login = tk.Button(root, text="Login", command=login, font=("Helvetica", 12))
button_signup = tk.Button(root, text="Sign Up", command=signup, font=("Helvetica", 12))

# Start with the login screen
show_login_screen()

root.mainloop()
