
import tkinter as tk # Importing the GUI toolkit for creating the application's interface
from tkinter import messagebox, filedialog  # For displaying messages and handling file dialog operations
import hashlib  # For hashing operations such as generating hashes
import os    # For interacting with the file system
import re # For handling regular expressions (used for validation purposes)
from Crypto.Cipher import AES  # For encryption and decryption using the AES algorithm
from Crypto.Util.Padding import pad  # For adding padding to plaintext
from Crypto.Util.Padding import unpad  # For removing padding from decrypted text
from Crypto.Random import get_random_bytes  # For generating random bytes for encryption keys
from docx import Document  # For handling Word document files (used for opening and editing .docx files)
from PyPDF2 import PdfReader, PdfWriter    # For reading and writing PDF files
from pptx import Presentation   # For handling PowerPoint presentations (used for opening and editing .pptx files)
import openpyxl  # For interacting with Excel files (used for opening and editing .xlsx files)
import tempfile   # For creating and managing temporary files and directories
from Crypto.PublicKey import RSA # For handling RSA encryption and decryption (used for securing file operations)
from Crypto.Cipher import PKCS1_OAEP  # For RSA encryption/decryption using OAEP padding
import socket # For socket programming (used for communication between client and server)


USER_DATA_FILE = "users.txt"  # file to store users credinationals
UPLOADS_DIR = "uploads"  # Directory to store uploaded files


# Ensure the uploads directory exists
if not os.path.exists(UPLOADS_DIR):
    os.makedirs(UPLOADS_DIR)

# Global variable to keep track of UI components
# This list stores references to the currently active widgets in the UI
current_widgets = []

# Function to clear all widgets in the current UI
# This ensures that any previous screen or components are removed before displaying new ones
def clear_ui():
    for widget in current_widgets:  # Iterate over all active widgets
        widget.pack_forget()  # Remove the widget from the UI
    current_widgets.clear()  # Clear the list to reset the UI state


# Function to generate a salt
def generate_salt():
    return os.urandom(16)  # Generates a random salt of 16 bytes

# Function to hash password with a salt
def hash_password(password, salt):
    return hashlib.sha256(salt + password.encode()).hexdigest()

# Function to save user data (username, hashed password, and salt)
def save_user(username, hashed_password, salt):
    with open(USER_DATA_FILE, "a") as file:
        file.write(f"{username},{hashed_password},{salt.hex()}\n")


def verify_user(username, password):
    if not os.path.exists(USER_DATA_FILE): # Check if the user data file exists. If not, return False as no users are registered 
        return False  # No users file yet

    with open(USER_DATA_FILE, "r") as file:    # Open the user data file in read mode
        for line in file:         # Loop through each line in the file
            line = line.strip()  # Remove any leading or  whitespace from the line
            if not line:  # Skip empty lines
                continue
            parts = line.split(",")
            
            if len(parts) != 3: # Split the line by commas. It is expected to have 3 parts: username, hashed password, and salt
                print(f"Skipping invalid line: {line}")  # Log invalid lines for debugging
                continue  # Skip lines that don't have exactly 3 parts
            
            
            # Assign the parts to meaningful variable names
            saved_username, saved_hashed_password, saved_salt = parts
            if username == saved_username:  # Check if the username from the file matches the input username
                salt = bytes.fromhex(saved_salt)  # Convert the salt back from hex
                hashed_input_password = hash_password(password, salt)   # Hash the input password with the retrieved salt
                if hashed_input_password == saved_hashed_password:  # Compare the hashed password with the saved hashed password
                    return True  # Return True if the passwords match

    return False   # If no matching username or password is found, return False




# Function to validate password complexity
def is_password_complex(password):
    if len(password) < 8:
        return False
    if not re.search(r"[A-Z]", password):  # At least one uppercase letter
        return False
    if not re.search(r"[a-z]", password):  # At least one lowercase letter
        return False
    if not re.search(r"\d", password):  # At least one digit
        return False
    if not re.search(r"[!@#$%^&*(),.?\":{}|<>_]", password):  # At least one special character
        return False
    return True

# Signup function
def signup():
    username = entry_username.get()
    password = entry_password.get()

    if not username or not password:
        messagebox.showwarning("Input Error", "Both fields are required!")
        return

    if not is_password_complex(password):
        messagebox.showwarning(
            "Weak Password",
            "Password must be at least 8 characters long, contain uppercase and lowercase letters, a number, and a special character.",
        )
        return

    salt = generate_salt()  # Generate a salt for the user
    hashed_password = hash_password(password, salt)  # Hash the password with the salt
    save_user(username, hashed_password, salt)
    messagebox.showinfo("Signup Success", "Account created successfully! Please log in now.")

    # Clear the username and password fields
    entry_username.delete(0, tk.END)
    entry_password.delete(0, tk.END)
    show_login_screen()

# Login function
def login():
    username = entry_username.get()
    password = entry_password.get()

    if not username or not password:
        messagebox.showwarning("Input Error", "Both fields are required!")
        return

    if verify_user(username, password):
        messagebox.showinfo("Login Success", f"Welcome, {username}!")
        show_upload_screen()  # Show the upload screen on successful login
    else:
        messagebox.showerror("Login Failed", "Invalid username or password! Please sign up if you don't have an account.")
        entry_password.delete(0, tk.END)


# Function to show the login screen
def show_login_screen():
    clear_ui()  # Clear any existing UI elements from the screen

    # Set the title text for the login screen and apply a font style
    label_title.config(text="File Manager - Please Sign Up or Log In", font=("Helvetica", 18))
    
    # Display the title label with some padding
    label_title.pack(pady=10)
    current_widgets.append(label_title)  # Keep track of the label in current_widgets

    # Display the username label
    label_username.pack()
    current_widgets.append(label_username)  # Track the label in current_widgets

    # Display the username input field with some padding
    entry_username.pack(pady=5)
    current_widgets.append(entry_username)  # Track the input field in current_widgets

    # Display the password label
    label_password.pack()
    current_widgets.append(label_password)  # Track the label in current_widgets

    # Display the password input field with some padding
    entry_password.pack(pady=5)
    current_widgets.append(entry_password)  # Track the input field in current_widgets

    # Display the login button with some padding
    button_login.pack(pady=10)
    current_widgets.append(button_login)  # Track the button in current_widgets

    # Display the signup button with some padding
    button_signup.pack(pady=10)
    current_widgets.append(button_signup)  # Track the button in current_widgets

    
# Function to show the file upload screen
def show_upload_screen():
    clear_ui()  # Clear any existing UI elements from the screen

    # Set the title text for the upload screen and apply a font style
    label_title.config(text="File Upload", font=("Helvetica", 18))
    
    # Display the title label with some padding
    label_title.pack(pady=10)
    current_widgets.append(label_title)  # Keep track of the label in current_widgets

    # Create and display a label with information about selecting a file
    label_info = tk.Label(root, text="Select a file to upload", font=("Helvetica", 12))
    label_info.pack(pady=10)
    current_widgets.append(label_info)  # Track the label in current_widgets

    # Create a button to browse and upload a file, linking it to the upload_file function
    button_browse = tk.Button(root, text="upload file", command=upload_file, font=("Helvetica", 12))
    button_browse.pack(pady=10)
    current_widgets.append(button_browse)  # Track the button in current_widgets

    # Create a logout button that takes the user back to the login screen
    button_logout = tk.Button(root, text="Logout", command=show_login_screen, font=("Helvetica", 12))
    button_logout.pack(pady=10)
    current_widgets.append(button_logout)  # Track the button in current_widgets

    
    

def encrypt_file(file_path, destination_path): 
    key = get_random_bytes(16)  # Generate a random 16-byte key for AES encryption (AES requires a 128-bit key)
    
    cipher = AES.new(key, AES.MODE_CBC)  # Create a new AES cipher object using CBC mode (Cipher Block Chaining)

    with open(file_path, 'rb') as file:  # Open the file to be encrypted in binary read mode
        file_data = file.read()  # Read the file data into memory
        
        # Encrypt the file data using the AES cipher, padding it to the required block size
        encrypted_data = cipher.encrypt(pad(file_data, AES.block_size))

    with open(destination_path, 'wb') as enc_file:  # Open the destination path to write the encrypted file in binary mode
        enc_file.write(cipher.iv)  # Write the initialization vector (IV) used for encryption
        enc_file.write(encrypted_data)  # Write the encrypted file data

    # Save the encryption key in a separate file for future decryption (never store it alongside the encrypted file)
    key_file_path = destination_path + ".key"  # Generate the path for the key file by appending '.key' to the destination file path
    with open(key_file_path, 'wb') as key_file:  # Open the key file in binary write mode
        key_file.write(key)  # Write the encryption key to the key file

    return key_file_path  # Return the path to the saved key file, which can be used for decryption


# Function to handle file upload
def upload_file():
    # Open a file dialog to select a file, filtering for specific file types
    file_path = filedialog.askopenfilename(title="Select a file", filetypes=(("Word Files", "*.docx"), ("PDF Files", "*.pdf"), ("PowerPoint Files", "*.pptx"), ("Excel Files", "*.xlsx"), ("All Files", "*.*")))
    
    # If the user selects a file (i.e., the file_path is not empty)
    if file_path:
        # Get the file name from the full path
        file_name = os.path.basename(file_path)
        
        # Define the destination path where the file will be uploaded
        destination = os.path.join(UPLOADS_DIR, file_name)

        # Check if the file already exists in the destination folder
        if os.path.exists(destination):
            # Show a warning if the file already exists
            messagebox.showwarning("File Already Uploaded", f"The file '{file_name}' has already been uploaded.")
            return  # Exit the function to prevent overwriting the file
        
        # Open the selected file in binary read mode
        with open(file_path, 'rb') as source_file:
            # Open the destination file in binary write mode
            with open(destination, 'wb') as dest_file:
                # Copy the content of the selected file to the destination
                dest_file.write(source_file.read())
                
        # Call the encryption function to encrypt the file after uploading
        encryption_key = encrypt_file(file_path, destination)
        
        # Show a success message after the file has been uploaded and encrypted
        messagebox.showinfo("File Upload", f"File '{file_name}' uploaded successfully!")



# Function to show uploaded files as clickable buttons
def show_uploaded_files():
    clear_ui()  # Clears the current UI elements to refresh the screen
    label_title.config(text="Uploaded Files", font=("Helvetica", 18))  # Set the title text and font for the screen
    label_title.pack(pady=10)  # Display the title label with padding
    current_widgets.append(label_title)  # Add the title label to the list of current UI widgets

    # Get the list of uploaded files (excluding .key files)
    uploaded_files = [file for file in os.listdir(UPLOADS_DIR) if not file.endswith(".key")]  # List all files in the UPLOADS_DIR, excluding .key files

    if not uploaded_files:  # Check if there are no uploaded files
        messagebox.showinfo("No Files", "No files have been uploaded yet.")  # Show a message box if no files are found
        show_upload_screen()  # Show the upload screen if there are no files
        return  # Exit the function if no files are available

    # Display uploaded files as clickable buttons
    for file in uploaded_files:  # Iterate over each uploaded file
        file_button = tk.Button(root, text=file, font=("Helvetica", 12), command=lambda file=file: file_options(file))  # Create a button for each file
        file_button.pack(pady=5)  # Display the file button with padding
        current_widgets.append(file_button)  # Add the button to the list of current UI widgets

    # Back button to go back to the upload screen
    button_back = tk.Button(root, text="Back", command=show_upload_screen, font=("Helvetica", 12))  # Create a back button
    button_back.pack(pady=10)  # Display the back button with padding
    current_widgets.append(button_back)  # Add the back button to the list of current UI widgets



    
# Function to handle the options for a selected file
def file_options(file_name):
    clear_ui()  # Clears the current UI elements to refresh the screen

    label_title.config(text=f"Options for {file_name}", font=("Helvetica", 18))  # Set the title to show which file options are being displayed
    label_title.pack(pady=10)  # Display the title label with padding
    current_widgets.append(label_title)  # Add the title label to the list of current UI widgets

    # Button to delete the file
    button_delete = tk.Button(root, text="Delete File", command=lambda: delete_file(file_name), font=("Helvetica", 12))  # Create a button for deleting the file
    button_delete.pack(pady=10)  # Display the delete button with padding
    current_widgets.append(button_delete)  # Add the delete button to the list of current UI widgets

    # Button to rename the file
    button_rename = tk.Button(root, text="Rename File", command=lambda: rename_file(file_name), font=("Helvetica", 12))  # Create a button for renaming the file
    button_rename.pack(pady=10)  # Display the rename button with padding
    current_widgets.append(button_rename)  # Add the rename button to the list of current UI widgets

    # Button to edit the file (added here)
    button_edit = tk.Button(root, text="Edit File", command=lambda: edit_file_in_window(file_name), font=("Helvetica", 12))  # Create a button for editing the file
    button_edit.pack(pady=10)  # Display the edit button with padding
    current_widgets.append(button_edit)  # Add the edit button to the list of current UI widgets

    # Button to decrypt and download the file
    button_download_decrypt = tk.Button(root, text="Download File", command=lambda: download_decrypt_file(file_name), font=("Helvetica", 12))  # Create a button for downloading and decrypting the file
    button_download_decrypt.pack(pady=10)  # Display the download button with padding
    current_widgets.append(button_download_decrypt)  # Add the download button to the list of current UI widgets

    # Back button to go back to the uploaded files list
    button_back = tk.Button(root, text="Back", command=show_uploaded_files, font=("Helvetica", 12))  # Create a back button to return to the uploaded files screen
    button_back.pack(pady=10)  # Display the back button with padding
    current_widgets.append(button_back)  # Add the back button to the list of current UI widgets



# Function to delete a selected file
def delete_file(file_name):
    # Construct the full file path by combining the uploads directory and the file name
    file_path = os.path.join(UPLOADS_DIR, file_name)
    
    try:
        # Try to remove the file at the specified path
        os.remove(file_path)
        
        # Show a success message when the file is deleted successfully
        messagebox.showinfo("File Deleted", f"The file '{file_name}' has been deleted successfully!")
        
        # Refresh the file list by displaying the updated list of uploaded files
        show_uploaded_files()
        
    except Exception as e:
        # If an error occurs, display an error message with the exception details
        messagebox.showerror("Error", f"An error occurred while deleting the file: {str(e)}")





# Function to rename a selected file
def rename_file(file_name):
    # Nested function to apply the renaming logic
    def apply_rename():
        # Get the new name entered by the user
        new_name = entry_new_name.get()
        
        # Check if the new name is empty and show a warning if true
        if not new_name:
            messagebox.showwarning("Invalid Name", "Please enter a valid name.")
            return
        
        # Create paths for the old and new file names and key files
        old_path = os.path.join(UPLOADS_DIR, file_name)
        new_path = os.path.join(UPLOADS_DIR, new_name)
        old_key_path = old_path + ".key"
        new_key_path = new_path + ".key"

        # Check if the new name already exists, and show a warning if true
        if os.path.exists(new_path):
            messagebox.showwarning("File Exists", f"A file with the name '{new_name}' already exists.")
            return

        try:
            # Rename the original file
            os.rename(old_path, new_path)
            # Rename the corresponding key file if it exists
            if os.path.exists(old_key_path):
                os.rename(old_key_path, new_key_path)
            # Show a success message and refresh the file list
            messagebox.showinfo("File Renamed", f"The file has been renamed to '{new_name}' successfully!")
            show_uploaded_files()  # Refresh the file list
        except Exception as e:
            # Show an error message if something goes wrong during the renaming
            messagebox.showerror("Error", f"An error occurred while renaming the file: {str(e)}")

    # Clear the current UI
    clear_ui()

    # Update the title with the current file's name
    label_title.config(text=f"Rename {file_name}", font=("Helvetica", 18))
    label_title.pack(pady=10)
    current_widgets.append(label_title)

    # Prompt the user to enter a new name
    label_new_name = tk.Label(root, text="Enter new file name:", font=("Helvetica", 12))
    label_new_name.pack(pady=10)
    current_widgets.append(label_new_name)

    # Create an entry field with the current file name as the default
    entry_new_name = tk.Entry(root, font=("Helvetica", 12))
    entry_new_name.pack(pady=10)
    entry_new_name.insert(0, file_name)  # Set the current file name as the default
    current_widgets.append(entry_new_name)

    # Create a button to apply the renaming action
    button_apply_rename = tk.Button(root, text="Apply Rename", command=apply_rename, font=("Helvetica", 12))
    button_apply_rename.pack(pady=10)
    current_widgets.append(button_apply_rename)

    # Create a cancel button to go back to the uploaded files list
    button_cancel = tk.Button(root, text="Cancel", command=show_uploaded_files, font=("Helvetica", 12))
    button_cancel.pack(pady=10)
    current_widgets.append(button_cancel)




    
def edit_file_in_window(file_name): # Define the function to edit a selected file
    clear_ui()   # Clear the current UI to set up the new screen
    label_title.config(text=f"Edit File: {file_name}", font=("Helvetica", 18))  # Set the title to indicate the file being edited
    label_title.pack(pady=10)  # Display the title label on the UI with some padding
    current_widgets.append(label_title)  # Add the title label to the list of current widgets

    file_path = os.path.join(UPLOADS_DIR, file_name)  # Create the full path to the file
    key_file_path = file_path + ".key"  # Path to the encryption key file

    # Create a temporary directory for decrypted files
    temp_dir = tempfile.mkdtemp()
    decrypted_file_path = os.path.join(temp_dir, file_name)

    # Decrypt the file before opening
    try:
        decrypt_file(file_path, key_file_path, decrypted_file_path)   # Call the decrypt function to get the decrypted file
    except Exception as e:
        messagebox.showerror("Error", f"Failed to decrypt file: {e}")
        show_uploaded_files()
        return

      # Attempt to load and display the content of the decrypted file based on its format
    try:
        if file_name.endswith(".docx"):  # If it's a Word document
            doc = Document(decrypted_file_path)  # Load the document
            content = "\n".join([para.text for para in doc.paragraphs])  # Extract text from paragraphs
        elif file_name.endswith(".pdf"):  # If it's a PDF document
            pdf_reader = PdfReader(decrypted_file_path)  # Load the PDF
            content = "\n".join([page.extract_text() for page in pdf_reader.pages])  # Extract text from each page
        elif file_name.endswith(".pptx"):  # If it's a PowerPoint presentation
            prs = Presentation(decrypted_file_path)  # Load the presentation
            content = "\n".join([slide.shapes.title.text if slide.shapes.title else '' for slide in prs.slides])  # Extract titles of each slide
        elif file_name.endswith(".xlsx"):  # If it's an Excel spreadsheet
            wb = openpyxl.load_workbook(decrypted_file_path)  # Load the workbook
            sheet = wb.active  # Get the active sheet
            content = "\n".join([str(cell.value) for row in sheet.iter_rows() for cell in row])  # Extract all cell values
        else:
            content = ""  # Default to empty if the file format is not supported
    except Exception as e:
        messagebox.showerror("Error", f"Failed to open file: {e}")  # Show error if loading the file fails
        show_uploaded_files()  # Return to the uploaded files screen
        return

    # Display the extracted content in a Text widget so the user can edit it
    text_edit = tk.Text(root, wrap=tk.WORD, height=15, width=40)  # Create a text widget with word wrapping
    text_edit.insert(tk.END, content)  # Insert the extracted content into the text widget
    text_edit.pack(pady=10)  # Pack the text widget into the UI with some padding
    current_widgets.append(text_edit)  # Add the text widget to the list of current widgets
    # Save edits back to the encrypted file
    def save_edits_and_exit():
        new_content = text_edit.get("1.0", tk.END).strip() # Get the content entered in the text widget and remove any extra whitespace

        try:
            # Check the file format and save the new content accordingly
            if file_name.endswith(".docx"):  # If the file is a Word document
                doc = Document()  # Create a new Document object
                doc.add_paragraph(new_content)  # Add the new content as a paragraph
                doc.save(decrypted_file_path)  # Save the modified document to the decrypted file path
            elif file_name.endswith(".pdf"):  # If the file is a PDF document
                pdf_writer = PdfWriter() # Create a new PDF writer object
                pdf_reader = PdfReader(decrypted_file_path) # Load the decrypted PDF file
                for page in pdf_reader.pages: # Add each page to the writer object
                    pdf_writer.add_page(page)
                # Modify content as needed
                with open(decrypted_file_path, 'wb') as f:
                    pdf_writer.write(f)
            elif file_name.endswith(".pptx"): # If the file is a PowerPoint presentation
                prs = Presentation()   # Create a new PowerPoint presentation
                slide = prs.slides.add_slide(prs.slide_layouts[0])   # Add a new slide
                if new_content:  # If there is new content
                    slide.shapes.title.text = new_content   # Set the slide's title to the new content
                prs.save(decrypted_file_path)  # Save the presentation to the decrypted file path
            elif file_name.endswith(".xlsx"): # If the file is an Excel spreadsheet
                wb = openpyxl.Workbook() # Create a new workbook object
                sheet = wb.active # Get the active sheet
                rows = new_content.splitlines()  # Split the new content into rows
                for i, row in enumerate(rows, start=1):  # Add each row to the spreadsheet
                    sheet[f"A{i}"] = row # Add content to the first column
                wb.save(decrypted_file_path)  # Save the workbook to the decrypted file path

            # Encrypt and save the updated file
            encrypt_file(decrypted_file_path, file_path)
            messagebox.showinfo("Success", f"File '{file_name}' has been updated!")
            root.destroy()  # Close the program after saving
        except Exception as e:
            messagebox.showerror("Error", f"Failed to save edits: {e}")

    # Add Save and Back buttons
    button_save = tk.Button(root, text="Save Edits", command=save_edits_and_exit, font=("Helvetica", 12))
    button_save.pack(pady=10)
    current_widgets.append(button_save)

    button_back = tk.Button(root, text="Back", command=show_uploaded_files, font=("Helvetica", 12))
    button_back.pack(pady=10)
    current_widgets.append(button_back)



def decrypt_file(encrypted_file_path, key_file_path, destination_path):
    # Check if the encrypted file exists
    if not os.path.exists(encrypted_file_path):
        raise FileNotFoundError(f"Encrypted file not found: {encrypted_file_path}")
    
    # Check if the key file exists
    if not os.path.exists(key_file_path):
        raise FileNotFoundError(f"Key file not found: {key_file_path}")

    try:
        # Open the encrypted file in binary read mode
        with open(encrypted_file_path, 'rb') as enc_file:
            # Read the initialization vector (IV) used for encryption (first 16 bytes)
            iv = enc_file.read(16)
            # Read the remaining encrypted data
            encrypted_data = enc_file.read()

        # Open the key file in binary read mode
        with open(key_file_path, 'rb') as key_file:
            # Read the encryption key
            key = key_file.read()

        # Initialize the AES cipher with the key, IV, and CBC mode
        cipher = AES.new(key, AES.MODE_CBC, iv)
        # Decrypt the encrypted data and remove padding
        decrypted_data = unpad(cipher.decrypt(encrypted_data), AES.block_size)

        # Open the destination file in binary write mode to save the decrypted data
        with open(destination_path, 'wb') as dec_file:
            dec_file.write(decrypted_data)

        # Return True to indicate the decryption was successful
        return True
    except Exception as e:
        # If any error occurs, raise a runtime error with a message
        raise RuntimeError(f"Decryption failed: {e}")



   


# Function to handle decryption and download of the file
def download_decrypt_file(file_name):
    encrypted_file_path = os.path.join(UPLOADS_DIR, file_name)
    key_file_path = encrypted_file_path + ".key"  # Assuming the key is stored with this extension
    
    # Check if the encrypted file and key file exist
    if not os.path.exists(encrypted_file_path):
        messagebox.showerror("Error", "Encrypted file not found.")
        return
    if not os.path.exists(key_file_path):
        messagebox.showerror("Error", "Key file not found.")
        return

    # Ask the user where to save the decrypted file
    decrypted_file_path = filedialog.asksaveasfilename(
        defaultextension=".dec", filetypes=[("All Files", "*.*")], title="Save Decrypted File"
    )

    if not decrypted_file_path:
        return  # User canceled the save dialog

    # Perform decryption
    if decrypt_file(encrypted_file_path, key_file_path, decrypted_file_path):
        messagebox.showinfo("Download Success", f"File successfully and saved as {decrypted_file_path}.")
    else:
        messagebox.showerror("Error", "Decryption failed.")

    # Ask user for a destination to save the decrypted file
    destination_path = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("DOCX Files", "*.docx"), ("All Files", "*.*")])

    if destination_path:
        success = decrypt_file(encrypted_file_path, key_file_path, destination_path)
        if success:
            messagebox.showinfo("The file is Successfully downloaed", f"File  and saved to '{destination_path}'.")
        else:
            messagebox.showerror("Decryption Failed", "Failed to decrypt the file.")



# Server address and port
SERVER_HOST = "127.0.0.1"
SERVER_PORT = 12345

# Function to load the server's public key
def load_server_public_key():
    try:
        # Open the public key file in binary read mode
        with open('server_public_key.pem', 'rb') as key_file:
            # Read the public key data and import it using the RSA module
            public_key = RSA.import_key(key_file.read())
        
        # Return the loaded public key
        return public_key
    except Exception as e:
        # If an error occurs (e.g., file not found, invalid key format), print the error message
        print(f"Error loading public key: {e}")
        # Return None if loading the public key fails
        return None



    
# Encrypt file using AES and the server's public RSA key
def encrypt_file_aes(file_path):
    try:
        # Generate a random AES key (256 bits) for encryption
        aes_key = get_random_bytes(32)  # AES key length is 256 bits (32 bytes)
        
        # Create a new AES cipher object in CBC (Cipher Block Chaining) mode using the generated key
        cipher = AES.new(aes_key, AES.MODE_CBC)

        # Open the file specified by the file_path in binary read mode
        with open(file_path, 'rb') as f:
            file_data = f.read()  # Read the entire file content into memory

        # Pad data to ensure its length is a multiple of 16 bytes (AES block size)
        padded_data = pad(file_data, AES.block_size)  # Use the 'pad' function to add padding
        encrypted_data = cipher.encrypt(padded_data)  # Encrypt the padded data using AES encryption

        # Generate a new file path to save the encrypted file (add .enc extension)
        encrypted_file_path = file_path + ".enc"
        
        # Open a new file in binary write mode to save the encrypted data
        with open(encrypted_file_path, 'wb') as f:
            f.write(cipher.iv)  # Save the Initialization Vector (IV) at the beginning of the file
            f.write(encrypted_data)  # Write the encrypted data after the IV

        # Load the server's public RSA key for encrypting the AES key
        server_public_key = load_server_public_key()
        if server_public_key is None:
            raise Exception("Failed to load server public key.")  # Raise an error if the public key is not loaded successfully

        # Create a new RSA cipher object using the loaded server public key
        rsa_cipher = PKCS1_OAEP.new(server_public_key)
        
        # Encrypt the AES key using the server's public RSA key
        encrypted_aes_key = rsa_cipher.encrypt(aes_key)

        # Return the encrypted file path and the encrypted AES key
        return encrypted_file_path, encrypted_aes_key
    except Exception as e:
        # Catch any exceptions that occur during encryption and print the error
        print(f"Error during encryption: {e}")
        return None, None  # Return None values if an error occurs
   # Function to upload and encrypt the file
   
   
   
def Upload_file():
    # Open a file dialog to allow the user to select a file for upload
    file_path = filedialog.askopenfilename(title="Select a file")
    
    # Check if the user selected a file and if the file exists
    if not file_path or not os.path.exists(file_path):
        # Show a warning if no file is selected or the file does not exist
        messagebox.showwarning("File Selection", "No file selected or file does not exist.")
        return  # Exit the function if the file is not valid

    # Call the function to encrypt the file using AES and the server's public RSA key
    encrypted_file_path, encrypted_aes_key = encrypt_file_aes(file_path)
    
    # Check if the file was successfully encrypted
    if encrypted_file_path and encrypted_aes_key:
        # Show a message confirming the file was shared successfully
        messagebox.showinfo("File shared", f"File shared successfully: {encrypted_file_path}")
        
        # Call the function to send the encrypted file and AES key to the server
        send_file_to_server(encrypted_file_path, encrypted_aes_key)
    else:
        # Show an error message if the file encryption failed
        messagebox.showerror("Encryption Error", "Failed to encrypt the file.")




# Function to send the encrypted file to the server
def send_file_to_server(encrypted_file_path, encrypted_aes_key):
    try:
        # Create a socket connection to the server
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            # Connect to the server using the defined host and port
            s.connect((SERVER_HOST, SERVER_PORT))

            # Send the file name length (4 bytes)
            file_name = os.path.basename(encrypted_file_path)  # Get the file name from the file path
            file_name_length = len(file_name).to_bytes(4, byteorder='big')  # Convert file name length to 4 bytes
            s.sendall(file_name_length)  # Send the length of the file name to the server

            # Send the file name
            s.sendall(file_name.encode('utf-8'))  # Send the file name in bytes (encoded)

            # Send the encrypted AES key (length first)
            s.sendall(len(encrypted_aes_key).to_bytes(4, byteorder='big'))  # Send the length of the encrypted AES key
            s.sendall(encrypted_aes_key)  # Send the encrypted AES key to the server

            # Send the encrypted file data
            with open(encrypted_file_path, 'rb') as file:
                # Read the file in chunks and send each chunk
                while True:
                    file_data = file.read(4096)  # Read the file in 4 KB chunks
                    if not file_data:
                        break  # Stop reading when there is no more data
                    s.sendall(file_data)  # Send the current chunk of file data to the server

            print("File sent to server successfully!")  # Print confirmation message when the file is sent
    except Exception as e:
        # Print an error message if something goes wrong during the file sending process
        print(f"Failed to send file: {str(e)}")




# Function to show file upload screen with options
def show_upload_screen():
    clear_ui()
    label_title.config(text="File Upload", font=("Helvetica", 18))
    label_title.pack(pady=10)
    current_widgets.append(label_title)

    # Option to upload a file
    button_upload = tk.Button(root, text="Upload a File", command=upload_file, font=("Helvetica", 12))
    button_upload.pack(pady=10)
    current_widgets.append(button_upload)
    
    
    button_share = tk.Button(root, text="Select a file to share",command=Upload_file ,font=("Helvetica", 12))
    button_share.pack(pady=10)
    current_widgets.append(button_share)

    # Option to show uploaded files
    button_show_files = tk.Button(root, text="Show Uploaded Files", command=show_uploaded_files, font=("Helvetica", 12))
    button_show_files.pack(pady=10)
    current_widgets.append(button_show_files)

    # Logout button
    button_logout = tk.Button(root, text="Logout", command=show_login_screen, font=("Helvetica", 12))
    button_logout.pack(pady=10)
    current_widgets.append(button_logout)

        
# Create the main window
root = tk.Tk()
root.title("File Manager")

# Create shared widgets
label_title = tk.Label(root, text="", font=("Helvetica", 18))

label_username = tk.Label(root, text="Username", font=("Helvetica", 12))
entry_username = tk.Entry(root, font=("Helvetica", 12))

label_password = tk.Label(root, text="Password", font=("Helvetica", 12))
entry_password = tk.Entry(root, font=("Helvetica", 12), show="*")

button_login = tk.Button(root, text="Login", command=login, font=("Helvetica", 12))
button_signup = tk.Button(root, text="Sign Up", command=signup, font=("Helvetica", 12))

# Start with the login screen
show_login_screen()

root.mainloop()
